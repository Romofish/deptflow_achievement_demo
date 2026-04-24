from __future__ import annotations

import io
from typing import Any, Iterable

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .ai_utils import build_rule_based_narrative
from .metrics_utils import calculate_report_metrics
from .slide_spec_utils import (
    STYLE_PRESETS,
    build_default_slide_spec,
    get_data_frame_for_slide,
    normalize_style_preset,
    sync_slide_spec_context,
)


TEMPLATE_VERSION = "deptflow-fixed-v3"

FIXED_TEMPLATE_SLIDES = [
    {"Slide": 1, "Name": "Title Page", "Content": "Report title, period, template metadata, and KPI cards"},
    {"Slide": 2, "Name": "Executive Summary", "Content": "Editable narrative plus calculated metric callouts"},
    {"Slide": 3, "Name": "Achievement Overview Dashboard", "Content": "Core KPIs, category mix, and impact mix"},
    {"Slide": 4, "Name": "Activity Type Breakdown", "Content": "Top activity types from filtered data"},
    {"Slide": 5, "Name": "Project / Study Highlights", "Content": "Study-level summary table"},
    {"Slide": 6, "Name": "People Contribution View", "Content": "Contributor summary from filtered data"},
    {"Slide": 7, "Name": "Quality / Complexity Highlights", "Content": "High-impact, quality, and CF indicators"},
    {"Slide": 8, "Name": "Appendix Detail Table", "Content": "Filtered achievement details"},
]

WHITE = RGBColor(255, 255, 255)
SOFT_LINE = RGBColor(218, 226, 235)


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _theme(slide_spec: dict[str, Any]) -> dict[str, RGBColor]:
    preset = STYLE_PRESETS[normalize_style_preset(slide_spec.get("style_preset"))]
    return {
        "accent": _rgb(preset["accent"]),
        "accent_2": _rgb(preset["accent_2"]),
        "accent_3": _rgb(preset["accent_3"]),
        "ink": _rgb(preset["ink"]),
        "muted": _rgb(preset["muted"]),
        "paper": _rgb(preset["paper"]),
        "surface": _rgb(preset["surface"]),
    }


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _set_line(shape, color: RGBColor, width: float = 0.75) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def _set_paragraph_font(paragraph, size: int, color: RGBColor, bold: bool = False) -> None:
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold


def _add_background(slide, colors: dict[str, RGBColor]) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    _set_fill(bg, colors["surface"])
    bg.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.14))
    _set_fill(band, colors["accent"])
    band.line.fill.background()


def _add_footer(
    slide,
    colors: dict[str, RGBColor],
    *,
    generated_at: str | None = None,
    ai_provider: str | None = None,
    style_preset: str = "atlas",
) -> None:
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.1), Inches(0.24))
    text = f"{TEMPLATE_VERSION} | Style: {style_preset}"
    if generated_at:
        text += f" | Generated: {generated_at}"
    if ai_provider:
        text += f" | AI: {ai_provider}"
    p = footer.text_frame.paragraphs[0]
    p.text = text
    _set_paragraph_font(p, 7, colors["muted"])


def _set_title(slide, colors: dict[str, RGBColor], title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(10.8), Inches(0.58))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    _set_paragraph_font(p, 24, colors["ink"], bold=True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.57), Inches(0.9), Inches(10.4), Inches(0.32))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        _set_paragraph_font(p2, 10, colors["muted"])


def _summary_cards(metrics: dict[str, Any]) -> list[tuple[str, Any]]:
    summary = metrics.get("summary", {})
    quality = metrics.get("quality_complexity", {})
    return [
        ("Achievements", summary.get("total_achievements", 0)),
        ("Studies / Workstreams", summary.get("unique_studies", 0)),
        ("Contributors", summary.get("unique_contributors", 0)),
        ("CF-related Count", summary.get("cf_total", 0)),
        ("High Impact", quality.get("high_impact_count", 0)),
        ("Quality / Inspection", quality.get("quality_inspection_count", 0)),
    ]


def _add_metric(slide, colors: dict[str, RGBColor], x: float, y: float, label: str, value: Any, accent: RGBColor | None = None) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.42), Inches(0.88))
    _set_fill(box, colors["paper"])
    _set_line(box, SOFT_LINE)
    box.text_frame.clear()
    p = box.text_frame.paragraphs[0]
    p.text = str(value)
    p.alignment = PP_ALIGN.CENTER
    _set_paragraph_font(p, 21, accent or colors["accent"], bold=True)
    p2 = box.text_frame.add_paragraph()
    p2.text = label
    p2.alignment = PP_ALIGN.CENTER
    _set_paragraph_font(p2, 7, colors["muted"])


def _add_metric_grid(slide, colors: dict[str, RGBColor], metrics: dict[str, Any], x: float, y: float, cols: int = 3) -> None:
    accents = [colors["accent"], colors["accent_2"], colors["accent_3"]]
    for idx, (label, value) in enumerate(_summary_cards(metrics)[:6]):
        row, col = divmod(idx, cols)
        _add_metric(slide, colors, x + col * 2.66, y + row * 1.05, label, value, accents[idx % len(accents)])


def _add_text_block(
    slide,
    colors: dict[str, RGBColor],
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_size: int = 13,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    paragraphs = [part.strip() for part in str(text).splitlines() if part.strip()] or [""]
    for idx, paragraph_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = paragraph_text
        _set_paragraph_font(p, font_size, colors["ink"])
        if idx:
            p.space_before = Pt(5)


def _add_bullets(slide, colors: dict[str, RGBColor], x: float, y: float, w: float, h: float, lines: Iterable[str], font_size: int = 11) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = str(line)
        p.level = 0
        _set_paragraph_font(p, font_size, colors["ink"])


def _add_chart(slide, colors: dict[str, RGBColor], x: float, y: float, w: float, h: float, title: str, data: pd.DataFrame, chart_type: str) -> None:
    if data.empty or data.shape[1] < 2 or chart_type == "none":
        _add_text_block(slide, colors, x, y, w, h, "No chart data available.", font_size=12)
        return

    chart_data = CategoryChartData()
    view = data.iloc[:, :2].copy()
    chart_data.categories = [str(value)[:28] for value in view.iloc[:, 0].tolist()]
    values = pd.to_numeric(view.iloc[:, 1], errors="coerce").fillna(0).astype(int).tolist()
    chart_data.add_series(title, values)

    chart_enum = XL_CHART_TYPE.COLUMN_CLUSTERED if chart_type == "column" else XL_CHART_TYPE.BAR_CLUSTERED
    chart = slide.shapes.add_chart(chart_enum, Inches(x), Inches(y), Inches(w), Inches(h), chart_data).chart
    chart.has_legend = False
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)


def _truncate(value: Any, max_chars: int = 110) -> str:
    if pd.isna(value):
        return ""
    text = str(value).replace("\n", " ").strip()
    return text if len(text) <= max_chars else text[: max_chars - 1] + "..."


def _add_table(slide, colors: dict[str, RGBColor], x: float, y: float, w: float, h: float, df: pd.DataFrame, max_rows: int = 8) -> None:
    view = df.head(max_rows).copy()
    if view.empty:
        view = pd.DataFrame({"Message": ["No records found for selected filters."]})

    rows, cols = view.shape[0] + 1, view.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table

    for col_idx, col in enumerate(view.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col)
        _set_fill(cell, colors["ink"])
        for paragraph in cell.text_frame.paragraphs:
            _set_paragraph_font(paragraph, 7, WHITE, bold=True)

    body_size = 5 if cols >= 7 else 6
    for row_idx, (_, row) in enumerate(view.iterrows(), start=1):
        for col_idx, col in enumerate(view.columns):
            cell = table.cell(row_idx, col_idx)
            cell.text = _truncate(row[col], max_chars=95 if cols < 7 else 58)
            for paragraph in cell.text_frame.paragraphs:
                _set_paragraph_font(paragraph, body_size, colors["ink"])


def _render_slide(
    prs: Presentation,
    slide_spec: dict[str, Any],
    slide_cfg: dict[str, Any],
    df: pd.DataFrame,
    metrics: dict[str, Any],
    *,
    generated_at: str | None,
    ai_provider: str,
) -> None:
    colors = _theme(slide_spec)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _add_background(slide, colors)
    _set_title(slide, colors, str(slide_cfg.get("title", "")), str(slide_cfg.get("subtitle", "")))

    kind = slide_cfg.get("kind", "table")
    chart_type = slide_cfg.get("chart_type", "none")
    layout = slide_cfg.get("layout_variant", "table_full")
    table_df = get_data_frame_for_slide(slide_cfg, df, metrics)
    top_n = int(slide_cfg.get("top_n") or 8)

    if kind == "title":
        _add_text_block(slide, colors, 0.72, 1.55, 5.8, 1.1, slide_cfg.get("narrative", ""), font_size=18)
        _add_text_block(slide, colors, 0.72, 2.55, 5.9, 0.75, "Fixed 8-slide template. Data and charts reflect current filters only.", font_size=13)
        _add_metric_grid(slide, colors, metrics, 0.75, 4.1, cols=3)
    elif kind == "narrative":
        _add_text_block(slide, colors, 0.72, 1.32, 8.0, 4.15, slide_cfg.get("narrative", ""), font_size=14)
        _add_metric_grid(slide, colors, metrics, 9.05, 1.36, cols=1)
        _add_bullets(
            slide,
            colors,
            0.85,
            5.75,
            11.3,
            0.75,
            ["Narrative is editable. Numeric values are generated by application metrics, not by AI."],
            font_size=9,
        )
    elif layout == "dashboard" or kind == "dashboard":
        _add_metric_grid(slide, colors, metrics, 0.65, 1.24, cols=3)
        _add_chart(slide, colors, 0.72, 3.05, 5.7, 3.45, "Primary Breakdown", table_df, chart_type)
        impact_df = pd.DataFrame(metrics.get("impact_breakdown", {}).items(), columns=["Impact", "Count"])
        _add_chart(slide, colors, 6.9, 3.05, 5.3, 3.45, "Impact Level Breakdown", impact_df, "column")
    elif layout == "chart_table" and chart_type != "none":
        _add_chart(slide, colors, 0.72, 1.4, 6.2, 5.45, str(slide_cfg.get("title", "")), table_df, chart_type)
        _add_table(slide, colors, 7.25, 1.48, 5.25, 4.85, table_df, max_rows=top_n)
    elif layout == "table_chart" and chart_type != "none":
        _add_table(slide, colors, 0.55, 1.42, 6.0, 5.2, table_df, max_rows=top_n)
        _add_chart(slide, colors, 6.95, 1.42, 5.55, 5.2, str(slide_cfg.get("title", "")), table_df, chart_type)
    elif layout == "split":
        _add_table(slide, colors, 0.55, 1.42, 7.0, 5.2, table_df, max_rows=top_n)
        if chart_type != "none":
            _add_chart(slide, colors, 7.95, 1.42, 4.55, 5.2, str(slide_cfg.get("title", "")), table_df, chart_type)
        else:
            _add_metric_grid(slide, colors, metrics, 8.05, 1.55, cols=1)
    elif layout == "cards_table" or kind == "quality":
        _add_metric_grid(slide, colors, metrics, 0.72, 1.22, cols=3)
        _add_table(slide, colors, 0.45, 3.15, 12.45, 3.75, table_df, max_rows=top_n)
    elif layout == "hero":
        _add_text_block(slide, colors, 0.72, 1.45, 4.5, 2.2, slide_cfg.get("narrative") or str(slide_cfg.get("data_source", "")), font_size=17)
        _add_table(slide, colors, 5.45, 1.42, 7.1, 5.2, table_df, max_rows=top_n)
    else:
        _add_table(slide, colors, 0.42, 1.3, 12.45, 5.75, table_df, max_rows=top_n)

    _add_footer(
        slide,
        colors,
        generated_at=generated_at,
        ai_provider=ai_provider,
        style_preset=str(slide_spec.get("style_preset", "atlas")),
    )


def generate_achievement_ppt(
    df: pd.DataFrame,
    report_title: str,
    period_label: str,
    narrative_text: str | None = None,
    ai_provider: str = "rule_based",
    ai_model: str = "",
    generated_at: str | None = None,
    filters: dict[str, Any] | None = None,
    metrics: dict[str, Any] | None = None,
    slide_spec: dict[str, Any] | None = None,
) -> bytes:
    """Generate editable PPTX from filtered achievement data and a fixed SlideSpec."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    metrics = metrics or calculate_report_metrics(df)
    default_narrative = narrative_text or build_rule_based_narrative(metrics, period_label)
    slide_spec = slide_spec or build_default_slide_spec(
        metrics,
        df,
        report_title=report_title,
        period_label=period_label,
        narrative_text=default_narrative,
    )
    slide_spec = sync_slide_spec_context(
        slide_spec,
        report_title=report_title,
        period_label=period_label,
        default_narrative=default_narrative,
    )
    if narrative_text:
        for slide_cfg in slide_spec.get("slides", []):
            if int(slide_cfg.get("id", 0)) == 2:
                slide_cfg["narrative"] = narrative_text

    props = prs.core_properties
    props.title = report_title
    props.subject = f"{TEMPLATE_VERSION}; ai_provider={ai_provider}; generated_at={generated_at or ''}"[:255]
    props.comments = (
        f"Trace stored in report_history.csv; style={slide_spec.get('style_preset', 'atlas')}; "
        f"rows={len(df)}; filters={bool(filters)}"
    )[:255]

    for slide_cfg in slide_spec.get("slides", []):
        _render_slide(prs, slide_spec, slide_cfg, df, metrics, generated_at=generated_at, ai_provider=ai_provider)

    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()
